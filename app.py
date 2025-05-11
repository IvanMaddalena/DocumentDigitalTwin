from flask import Flask, render_template, request, jsonify, redirect, url_for
import os
import PyPDF2
import markdown2
from pptx import Presentation
from docx import Document
import time
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings.huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_community.llms import HuggingFaceEndpoint
from langchain.chains import RetrievalQA
from flask_sqlalchemy import SQLAlchemy

app = Flask(__name__)
UPLOAD_FOLDER = 'uploads'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///rooms.db'  # Il percorso del database SQLite
db = SQLAlchemy(app)

pdf_data = {}
rooms = {}


class Room(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(100), unique=True, nullable=False)
    password = db.Column(db.String(100), nullable=False)


class Documenti(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    filename = db.Column(db.String(255), unique=True, nullable=False)
    content = db.Column(db.Text, nullable=False)  # Campo per il contenuto del documento
    room_id = db.Column(db.Integer, db.ForeignKey('room.id'), nullable=False)
    room = db.relationship('Room', backref=db.backref('documents', lazy=True))


with app.app_context():
    db.create_all()


def extract_text_from_pdf(file_path):
    text = ''
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        for page_num in range(len(reader.pages)):
            text += reader.pages[page_num].extract_text()
    return text


# Funzione per estrarre il testo da un file Markdown (.md)
def extract_text_from_md(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        markdown_text = file.read()
        html = markdown2.markdown(markdown_text)
        # Rimuovi i tag HTML per ottenere solo il testo
        text = html
    return text


# Funzione per estrarre il testo da un file di testo (.txt)
def extract_text_from_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()


# Funzione per estrarre il testo da un file PowerPoint (.pptx)
def extract_text_from_pptx(file_path):
    presentation = Presentation(file_path)
    text = ''
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text
    return text


# Funzione per estrarre il testo da un file Word (.docx)
def extract_text_from_docx(file_path):
    doc = Document(file_path)
    full_text = []
    for paragraph in doc.paragraphs:
        full_text.append(paragraph.text)
    return '\n'.join(full_text)


# Funzione per creare una nuova room
def create_room(room_name, password):
    new_room = Room(name=room_name, password=password)
    db.session.add(new_room)
    db.session.commit()


# Funzione per verificare se una room esiste
def room_exists(room_name):
    return Room.query.filter_by(name=room_name).first() is not None


# Funzione per verificare la password di una room
def check_password(room_name, password):
    room = Room.query.filter_by(name=room_name).first()
    return room is not None and room.password == password


def process_question(pdf_title, question):
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], pdf_title)
    if pdf_title.endswith('.txt'):
        text = extract_text_from_txt(file_path)
    elif pdf_title.endswith('.pptx'):
        text = extract_text_from_pptx(file_path)
    elif pdf_title.endswith('.pdf'):
        text = extract_text_from_pdf(file_path)
    elif pdf_title.endswith('.md'):
        text = extract_text_from_md(file_path)
    elif pdf_title.endswith('.docx'):
        text = extract_text_from_docx(file_path)
    else:
        text = "Placeholder text for non-supported formats"

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=200, chunk_overlap=20)
    chunks = text_splitter.split_text(text)

    embeddings = HuggingFaceEmbeddings()
    vectorStore = FAISS.from_texts(chunks, embeddings)

    os.environ["HUGGINGFACEHUB_API_TOKEN"] = "INSERT YOUR TOKEN"
    llm = HuggingFaceEndpoint(repo_id="tiiuae/falcon-7b-instruct", temperature=0.1,
                              model_kwargs={"max_length": 512})

    chain = RetrievalQA.from_chain_type(llm=llm, chain_type="stuff", retriever=vectorStore.as_retriever())

    answer = chain.invoke(question)
    return answer['result']


# Funzione per processare tutte le domande e ottenere le risposte
def process_all_questions(pdf_data):
    for pdf_title, data in pdf_data.items():
        question = data['question']
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], pdf_title)

        if pdf_title.endswith('.txt'):
            text = extract_text_from_txt(file_path)
        elif pdf_title.endswith('.pptx'):
            text = extract_text_from_pptx(file_path)
        elif pdf_title.endswith('.pdf'):
            text = extract_text_from_pdf(file_path)
        elif pdf_title.endswith('.md'):
            text = extract_text_from_md(file_path)
        elif pdf_title.endswith('.docx'):
            text = extract_text_from_docx(file_path)
        else:
            text = "Placeholder text for non-supported formats"

        text_splitter = RecursiveCharacterTextSplitter(chunk_size=200, chunk_overlap=20)
        chunks = text_splitter.split_text(text)

        embeddings = HuggingFaceEmbeddings()
        vectorStore = FAISS.from_texts(chunks, embeddings)

        os.environ["HUGGINGFACEHUB_API_TOKEN"] = "INSERT YOUR TOKEN"
        llm = HuggingFaceEndpoint(repo_id="tiiuae/falcon-7b-instruct", temperature=0.1,
                                  model_kwargs={"max_length": 512})

        chain = RetrievalQA.from_chain_type(llm=llm, chain_type="stuff", retriever=vectorStore.as_retriever())

        answer = chain.invoke(question)
        data['answer'] = answer['result']


def generate_unique_filename(filename):
    base_name, extension = os.path.splitext(filename)
    counter = 1
    while True:
        new_filename = f"{base_name}_{counter}{extension}"
        if not Documenti.query.filter_by(filename=new_filename).first():
            # Il nome del file è unico
            return new_filename
        counter += 1


# Route per la pagina della room
@app.route("/room/<room_name>", methods=["GET", "POST"])
def room(room_name):
    room = Room.query.filter_by(name=room_name).first()
    # Recupera tutti i documenti associati alla stanza
    documents = Documenti.query.filter_by(room=room).all()
    if room is None:
        return render_template('room_not_found.html', room_name=room_name)

    if request.method == "POST":
        if 'all_pdfs' in request.form:
            global_question = request.form['global_question']
            for pdf_title in pdf_data.keys():
                pdf_data[pdf_title]['question'] = global_question
            return render_template('room.html', room=room, pdf_data=pdf_data)
        elif 'send_all_questions' in request.form:
            process_all_questions(pdf_data)  # Invia tutte le domande al sistema NLP
            return render_template('room.html', room=room, pdf_data=pdf_data)
        elif 'remove_document_id' in request.form:
            document_id_to_remove = int(request.form['remove_document_id'])
            document_to_remove = Documenti.query.get(document_id_to_remove)
            if document_to_remove:
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], document_to_remove.filename)
                if os.path.exists(file_path):
                    # Rimuovi il documento dal database
                    db.session.delete(document_to_remove)
                    db.session.commit()
                    # Rimuovi anche dalla lista dei documenti caricati in memoria
                    pdf_data.pop(document_to_remove.filename, None)
                    # Rimuovi il file dal sistema di file
                    os.remove(file_path)
                else:
                    db.session.delete(document_to_remove)
                    db.session.commit()
                    pdf_data.pop(document_to_remove.filename, None)
                    print("Il file non esiste nel percorso specificato:", file_path)
            # Redirect alla stessa pagina dopo la rimozione
            return redirect(url_for('room', room_name=room_name))
        elif 'remove_pdf' in request.form:
            pdf_title_to_remove = request.form['remove_pdf']
            if pdf_title_to_remove in pdf_data:
                del pdf_data[pdf_title_to_remove]
                os.remove(os.path.join(app.config['UPLOAD_FOLDER'], pdf_title_to_remove))
            return render_template('room.html', room=room, pdf_data=pdf_data)
        else:
            if 'files[]' in request.files:
                files = request.files.getlist('files[]')
                for uploaded_file in files:
                    filename = uploaded_file.filename
                    filename = generate_unique_filename(filename)
                    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                    uploaded_file.save(file_path)
                    # Estrai il contenuto del file
                    if filename.endswith('.txt'):
                        content = extract_text_from_txt(file_path)
                    elif filename.endswith('.pptx'):
                        content = extract_text_from_pptx(file_path)
                    elif filename.endswith('.pdf'):
                        content = extract_text_from_pdf(file_path)
                    elif filename.endswith('.md'):
                        content = extract_text_from_md(file_path)
                    elif filename.endswith('.docx'):
                        content = extract_text_from_docx(file_path)
                    else:
                        content = "Placeholder text for non-supported formats"
                    new_document = Documenti(filename=filename, content=content, room=room)
                    db.session.add(new_document)
                    db.session.commit()
                    pdf_data[filename] = {"pdf_title": filename, "question": None, "answer": None}
                pdf_titles = request.form.getlist('pdf_title')
                questions = request.form.getlist('question')

                for pdf_title, question in zip(pdf_titles, questions):
                    pdf_data[pdf_title]['question'] = question

                    file_path = os.path.join(app.config['UPLOAD_FOLDER'], pdf_title)
                    if pdf_title.endswith('.txt'):
                        text = extract_text_from_txt(file_path)
                    elif pdf_title.endswith('.pptx'):
                        text = extract_text_from_pptx(file_path)
                    elif pdf_title.endswith('.pdf'):
                        text = extract_text_from_pdf(file_path)
                    elif pdf_title.endswith('.md'):
                        text = extract_text_from_md(file_path)
                    elif pdf_title.endswith('.docx'):
                        text = extract_text_from_docx(file_path)
                    else:
                        text = "Placeholder text for non-supported formats"

                    text_splitter = RecursiveCharacterTextSplitter(chunk_size=200, chunk_overlap=20)
                    chunks = text_splitter.split_text(text)

                    embeddings = HuggingFaceEmbeddings()
                    vectorStore = FAISS.from_texts(chunks, embeddings)

                    os.environ["HUGGINGFACEHUB_API_TOKEN"] = "INSERT YOUR TOKEN"
                    llm = HuggingFaceEndpoint(repo_id="tiiuae/falcon-7b-instruct", temperature=0.1,
                                              model_kwargs={"max_length": 512})

                    chain = RetrievalQA.from_chain_type(llm=llm, chain_type="stuff",
                                                        retriever=vectorStore.as_retriever())

                    answer = chain.invoke(question)
                    pdf_data[pdf_title]['answer'] = answer['result']

                return render_template('room.html', room=room, pdf_data=pdf_data, documents=documents)

            elif 'load_document_id' in request.form:
                document_id = int(request.form['load_document_id'])
                document = Documenti.query.get(document_id)
                if document:
                    # Carica il documento dal database e aggiungilo a pdf_data se non è già presente
                    if document.filename not in pdf_data:
                        pdf_data[document.filename] = {"pdf_title": document.filename, "question": None, "answer": None}
                # Esegui altre operazioni se necessario
            return render_template('room.html', room=room, pdf_data=pdf_data, documents=documents)
    else:
        return render_template('room.html', room=room, pdf_data=pdf_data, documents=documents)


@app.route("/direct_access", methods=["POST"])
def direct_access():
    room_name = request.form.get('room_name')
    password = request.form.get('password')
    if room_exists(room_name) and check_password(room_name, password):
        # Redirect alla pagina della stanza
        return redirect(url_for('room', room_name=room_name))
    else:
        # Se la password non è corretta, rimani sulla stessa pagina con un messaggio di errore
        error_message = "Invalid room name or password"
        all_rooms = Room.query.all()
        return render_template('index.html', pdf_data=pdf_data, rooms=all_rooms, error_message=error_message)


# Endpoint per gestire il reindirizzamento alla pagina di inserimento della password per le stanze esistenti
@app.route("/join_room/<room_name>", methods=["GET"])
def join_room(room_name):
    return render_template('room_password.html', room_name=room_name)


@app.route("/", methods=["GET", "POST"])
def index():
    all_rooms = Room.query.all()
    search_query = request.args.get('search_query', '')
    if search_query:
        filtered_rooms = [room for room in all_rooms if search_query.lower() in room.name.lower()]
    else:
        filtered_rooms = all_rooms
    visible_rooms = filtered_rooms  # Visualizza solo le prime tre stanze
    # Se la query di ricerca è vuota, visualizza l'elenco originale limitato a tre stanze
    if not search_query:
        visible_rooms = all_rooms
    if request.method == "POST":
        if 'create_room' in request.form:
            room_name = request.form['room_name']
            password = request.form['password']
            # Verifica se il nome della stanza è già utilizzato
            if Room.query.filter_by(name=room_name).first():
                create_room_error = "Nome già utilizzato. Scegliere un altro nome."
                return render_template('index.html', pdf_data=pdf_data, rooms=visible_rooms,
                                       search_query=search_query, create_room_error=create_room_error)
            else:
                create_room(room_name, password)
                # Dopo aver creato la stanza, reindirizza l'utente alla pagina della stanza
                return redirect(url_for('room', room_name=room_name))
        elif 'join_room' in request.form:
            room_name = request.form['room_name']
            password = request.form['password']
            if room_exists(room_name) and check_password(room_name, password):
                # Redirect alla pagina della room
                return redirect(url_for('room', room_name=room_name))
            else:
                direct_access_error = "Invalid room name or password. Try again."
                return render_template('index.html', pdf_data=pdf_data, rooms=rooms,
                                       direct_access_error=direct_access_error)
        elif 'all_pdfs' in request.form:
            global_question = request.form['global_question']
            for pdf_title in pdf_data.keys():
                pdf_data[pdf_title]['question'] = global_question
            return render_template('index.html', pdf_data=pdf_data)
        elif 'send_all_questions' in request.form:
            process_all_questions(pdf_data)  # Invia tutte le domande al sistema NLP
            return render_template('index.html', pdf_data=pdf_data)
        elif 'remove_pdf' in request.form:
            pdf_title_to_remove = request.form['remove_pdf']
            if pdf_title_to_remove in pdf_data:
                del pdf_data[pdf_title_to_remove]
                os.remove(os.path.join(app.config['UPLOAD_FOLDER'], pdf_title_to_remove))
            return render_template('index.html', pdf_data=pdf_data)
        else:
            files = request.files.getlist('files[]')
            for uploaded_file in files:
                filename = uploaded_file.filename
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                uploaded_file.save(file_path)
                pdf_data[filename] = {"pdf_title": filename, "question": None, "answer": None}
            pdf_titles = request.form.getlist('pdf_title')
            questions = request.form.getlist('question')

            for pdf_title, question in zip(pdf_titles, questions):
                pdf_data[pdf_title]['question'] = question

                start_time = time.time()  # Salva il tempo di inizio dell'elaborazione
                first_token_time = None  # Variabile per salvare il tempo del primo token

                file_path = os.path.join(app.config['UPLOAD_FOLDER'], pdf_title)
                if pdf_title.endswith('.txt'):
                    text = extract_text_from_txt(file_path)
                elif pdf_title.endswith('.pptx'):
                    text = extract_text_from_pptx(file_path)
                elif pdf_title.endswith('.pdf'):
                    text = extract_text_from_pdf(file_path)
                elif pdf_title.endswith('.md'):
                    text = extract_text_from_md(file_path)
                elif pdf_title.endswith('.docx'):
                    text = extract_text_from_docx(file_path)
                else:
                    text = "Placeholder text for non-supported formats"

                text_splitter = RecursiveCharacterTextSplitter(chunk_size=200, chunk_overlap=20)
                chunks = text_splitter.split_text(text)

                embeddings = HuggingFaceEmbeddings()
                vectorStore = FAISS.from_texts(chunks, embeddings)

                os.environ["HUGGINGFACEHUB_API_TOKEN"] = "INSERT YOUR TOKEN"
                llm = HuggingFaceEndpoint(repo_id="tiiuae/falcon-7b-instruct", temperature=0.1,
                                          model_kwargs={"max_length": 512})

                chain = RetrievalQA.from_chain_type(llm=llm, chain_type="stuff", retriever=vectorStore.as_retriever())

                # Funzione callback per catturare il tempo del primo token
                def callback(token):
                    nonlocal first_token_time
                    if first_token_time is None:
                        first_token_time = time.time()

                answer = chain.invoke(question, callbacks=[callback])
                end_time = time.time()  # Salva il tempo di fine dell'elaborazione
                # Calcola i tempi di elaborazione
                processing_time = end_time - start_time  # Tempo totale di elaborazione
                time_to_first_token = first_token_time - start_time if first_token_time else None
                gen_time = end_time - (first_token_time if first_token_time else start_time)
                tokens_generated = len(answer['result'].split())
                speed = tokens_generated / gen_time if gen_time else None
                # Stampa i tempi e le metriche
                print(f"Tempo per il primo token: {time_to_first_token} secondi")
                print(f"Tempo di generazione totale: {gen_time} secondi")
                print(f"Velocità di generazione: {speed} token/secondo")
                print(f"Tempo di elaborazione per la domanda '{question}': {processing_time} secondi")
                pdf_data[pdf_title]['answer'] = answer['result']

            return render_template('index.html', pdf_data=pdf_data)
    else:
        return render_template('index.html', pdf_data=pdf_data, rooms=visible_rooms, search_query=search_query)


if __name__ == "__main__":
    app.run(debug=True)
